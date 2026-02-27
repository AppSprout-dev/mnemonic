#!/usr/bin/env python3
"""Generate Mnemonic sales pitch PPTX for Jason Bennitt — numbers-heavy version."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Design Tokens ──────────────────────────────────────────────────────────────
BG_NAVY      = RGBColor(0x1B, 0x2A, 0x4A)
BG_DARK      = RGBColor(0x12, 0x1E, 0x36)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY   = RGBColor(0xCC, 0xD6, 0xE0)
ACCENT_BLUE  = RGBColor(0x4F, 0xC3, 0xF7)
ACCENT_GREEN = RGBColor(0x66, 0xBB, 0x6A)
ACCENT_AMBER = RGBColor(0xFF, 0xCA, 0x28)
ACCENT_RED   = RGBColor(0xEF, 0x53, 0x50)
ACCENT_PURPLE= RGBColor(0xBA, 0x68, 0xC8)
MUTED_BLUE   = RGBColor(0x37, 0x87, 0xAF)
SHAPE_FILL   = RGBColor(0x25, 0x3B, 0x5E)
SHAPE_BORDER = RGBColor(0x4F, 0xC3, 0xF7)
DIM_TEXT      = RGBColor(0x8A, 0x9B, 0xB0)
DARK_GREEN_BG = RGBColor(0x1E, 0x3A, 0x2F)
DARK_RED_BG   = RGBColor(0x3A, 0x1E, 0x1E)
DARK_PURPLE_BG= RGBColor(0x2E, 0x1E, 0x3A)
DARK_AMBER_BG = RGBColor(0x3A, 0x2E, 0x1E)
CHECK_GREEN   = RGBColor(0x4C, 0xAF, 0x50)
X_RED         = RGBColor(0xE5, 0x73, 0x73)

FONT_NAME    = "Calibri"
MONO_FONT    = "Courier New"
SLIDE_W      = Inches(13.333)  # 16:9 widescreen
SLIDE_H      = Inches(7.5)

OUTPUT_PATH  = os.path.join(os.path.dirname(os.path.dirname(__file__)), "Mnemonic_Pitch.pptx")


# ── Helper Functions ───────────────────────────────────────────────────────────

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height):
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    return tf


def set_para(para, text, size=20, bold=False, color=WHITE, align=PP_ALIGN.LEFT, font_name=FONT_NAME):
    para.text = text
    para.alignment = align
    font = para.font
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    return para


def add_rounded_rect(slide, left, top, width, height, text, fill_color=SHAPE_FILL,
                     text_color=WHITE, font_size=12, border_color=SHAPE_BORDER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.name = FONT_NAME
    run.font.size = Pt(font_size)
    run.font.color.rgb = text_color
    run.font.bold = True
    return shape


def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    return shape


def add_slide_title(slide, title_text, subtitle_text=None):
    tf = add_textbox(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    set_para(tf.paragraphs[0], title_text, size=36, bold=True, color=WHITE)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.2), Inches(2.5), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()
    if subtitle_text:
        tf2 = add_textbox(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(0.5))
        set_para(tf2.paragraphs[0], subtitle_text, size=16, color=DIM_TEXT)


def add_stat_box(slide, x, y, w, h, number, label, num_color=WHITE, border=SHAPE_BORDER):
    box = add_rounded_rect(slide, x, y, w, h, "", font_size=14, border_color=border)
    box_tf = box.text_frame
    box_tf.clear()
    box_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    num_p = box_tf.paragraphs[0]
    num_p.text = number
    num_p.alignment = PP_ALIGN.CENTER
    num_p.font.name = FONT_NAME
    num_p.font.size = Pt(28)
    num_p.font.bold = True
    num_p.font.color.rgb = num_color
    label_p = box_tf.add_paragraph()
    label_p.text = label
    label_p.alignment = PP_ALIGN.CENTER
    label_p.font.name = FONT_NAME
    label_p.font.size = Pt(12)
    label_p.font.color.rgb = LIGHT_GRAY
    return box


def add_titled_bullets(tf, items, title_size=20, desc_size=15, title_color=ACCENT_BLUE, desc_color=LIGHT_GRAY):
    """Add title+description bullet pairs to a text frame."""
    for i, (title, desc) in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = title
        p.font.name = FONT_NAME
        p.font.size = Pt(title_size)
        p.font.bold = True
        p.font.color.rgb = title_color
        p.space_before = Pt(8)
        p.space_after = Pt(2)
        dp = tf.add_paragraph()
        dp.text = desc
        dp.font.name = FONT_NAME
        dp.font.size = Pt(desc_size)
        dp.font.color.rgb = desc_color
        dp.space_after = Pt(6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ═══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Slide 1 — Title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    tf = add_textbox(slide, Inches(1), Inches(1.6), Inches(11.3), Inches(1.5))
    set_para(tf.paragraphs[0], "MNEMONIC", size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    tf2 = add_textbox(slide, Inches(1), Inches(3.1), Inches(11.3), Inches(0.7))
    set_para(tf2.paragraphs[0], "The Memory Layer for AI Agents", size=28, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(3.9), Inches(3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    tf3 = add_textbox(slide, Inches(1), Inches(4.3), Inches(11.3), Inches(0.7))
    set_para(tf3.paragraphs[0], "Local-first  \u2022  Neuroscience-inspired  \u2022  Fully autonomous", size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

    # Market hook
    tf4 = add_textbox(slide, Inches(1), Inches(5.2), Inches(11.3), Inches(0.6))
    set_para(tf4.paragraphs[0], "$15.7B market  \u2022  42% CAGR  \u2022  1.3B AI agents by 2028", size=16, bold=True, color=ACCENT_AMBER, align=PP_ALIGN.CENTER)

    tf5 = add_textbox(slide, Inches(1), Inches(6.3), Inches(11.3), Inches(0.5))
    set_para(tf5.paragraphs[0], "Prepared for Jason Bennitt  \u2022  February 2026", size=14, color=DIM_TEXT, align=PP_ALIGN.CENTER)


def slide_02_problem(prs):
    """Slide 2 — The Problem + Market Stats."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "AI Has Amnesia")

    # Left: problem bullets
    tf = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(7.5), Inches(5))
    tf.paragraphs[0].space_after = Pt(0)
    bullets = [
        ("Every session starts from zero", "No memory of past decisions, errors, or insights. Context vanishes the moment you close the chat."),
        ("Knowledge workers lose 3+ hours/week", "Switching between tools, sessions, and projects means re-explaining the same things over and over."),
        ("Current solutions miss the mark", "Cloud-locked platforms are privacy-hostile. Vector databases offer search, not reasoning. Nothing truly learns."),
        ("AI agents can't learn from experience", "Without persistent memory, AI makes the same mistakes, asks the same questions, and never improves."),
    ]
    add_titled_bullets(tf, bullets, title_size=18, desc_size=14)

    # Right: market stats boxes
    add_stat_box(slide, Inches(9.0), Inches(1.8), Inches(3.5), Inches(1.3), "$15.7B", "AI Dev Tools Market by 2033", num_color=ACCENT_BLUE, border=ACCENT_BLUE)
    add_stat_box(slide, Inches(9.0), Inches(3.3), Inches(3.5), Inches(1.3), "42.3%", "CAGR Growth Rate", num_color=ACCENT_GREEN, border=ACCENT_GREEN)
    add_stat_box(slide, Inches(9.0), Inches(4.8), Inches(3.5), Inches(1.3), "1.3B", "AI Agents by 2028", num_color=ACCENT_AMBER, border=ACCENT_AMBER)


def slide_03_market(prs):
    """Slide 3 — NEW: Market Opportunity (TAM/SAM/SOM)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "A $15.7 Billion Opportunity")

    # TAM/SAM/SOM funnel — three nested boxes (largest on left)
    # TAM box
    add_rounded_rect(slide, Inches(0.8), Inches(1.8), Inches(3.5), Inches(4.5), "", fill_color=RGBColor(0x1A, 0x35, 0x55), border_color=ACCENT_BLUE)
    tf_tam = add_textbox(slide, Inches(0.9), Inches(1.9), Inches(3.3), Inches(4.3))
    p = tf_tam.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "TAM"
    r.font.name = FONT_NAME
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = ACCENT_BLUE
    p2 = tf_tam.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "$15.7B"
    r2.font.name = FONT_NAME
    r2.font.size = Pt(36)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    p3 = tf_tam.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(8)
    r3 = p3.add_run()
    r3.text = "AI in Software Development\nby 2033 (42.3% CAGR)"
    r3.font.name = FONT_NAME
    r3.font.size = Pt(12)
    r3.font.color.rgb = LIGHT_GRAY

    # SAM box
    add_rounded_rect(slide, Inches(4.6), Inches(1.8), Inches(3.5), Inches(4.5), "", fill_color=RGBColor(0x1A, 0x40, 0x35), border_color=ACCENT_GREEN)
    tf_sam = add_textbox(slide, Inches(4.7), Inches(1.9), Inches(3.3), Inches(4.3))
    p = tf_sam.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "SAM"
    r.font.name = FONT_NAME
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = ACCENT_GREEN
    p2 = tf_sam.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "$3.2B"
    r2.font.name = FONT_NAME
    r2.font.size = Pt(36)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    p3 = tf_sam.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(8)
    r3 = p3.add_run()
    r3.text = "AI Agent Memory &\nContext Management"
    r3.font.name = FONT_NAME
    r3.font.size = Pt(12)
    r3.font.color.rgb = LIGHT_GRAY

    # SOM box
    add_rounded_rect(slide, Inches(8.4), Inches(1.8), Inches(4.2), Inches(4.5), "", fill_color=RGBColor(0x40, 0x35, 0x1A), border_color=ACCENT_AMBER)
    tf_som = add_textbox(slide, Inches(8.5), Inches(1.9), Inches(4.0), Inches(4.3))
    p = tf_som.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "SOM (3-YEAR TARGET)"
    r.font.name = FONT_NAME
    r.font.size = Pt(14)
    r.font.bold = True
    r.font.color.rgb = ACCENT_AMBER
    p2 = tf_som.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text = "$12M ARR"
    r2.font.name = FONT_NAME
    r2.font.size = Pt(36)
    r2.font.bold = True
    r2.font.color.rgb = WHITE
    p3 = tf_som.add_paragraph()
    p3.alignment = PP_ALIGN.CENTER
    p3.space_after = Pt(12)
    r3 = p3.add_run()
    r3.text = "5,000 paid users +\n80 enterprise contracts"
    r3.font.name = FONT_NAME
    r3.font.size = Pt(12)
    r3.font.color.rgb = LIGHT_GRAY

    # Validation callouts at bottom
    tf_val = add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.6))
    set_para(tf_val.paragraphs[0],
             "16,000+ MCP servers  \u2022  97M SDK downloads/mo  \u2022  MCP donated to Linux Foundation  \u2022  53% of enterprises already using AI agents",
             size=13, color=DIM_TEXT, align=PP_ALIGN.CENTER)


def slide_04_solution(prs):
    """Slide 4 — The Solution."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "Meet Mnemonic")

    tf = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(4.5))
    tf.paragraphs[0].space_after = Pt(0)

    items = [
        "\u2713  Watches what you do and remembers what matters \u2014 autonomously",
        "\u2713  Runs entirely on your machine \u2014 no cloud, no subscriptions, no data leaks",
        "\u2713  AI agents get persistent memory across sessions via 10 MCP tools",
        "\u2713  Not just storage \u2014 it thinks: encodes, consolidates, dreams, discovers patterns",
        "\u2713  Learns from your feedback and self-improves over time",
    ]

    for i, text in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.name = FONT_NAME
        p.font.size = Pt(20)
        p.font.color.rgb = WHITE
        p.space_after = Pt(14)
        p.space_before = Pt(4)

    highlight = add_textbox(slide, Inches(1.5), Inches(5.8), Inches(10.3), Inches(0.6))
    set_para(highlight.paragraphs[0],
             '"Like giving your AI a brain that remembers, forgets wisely, and learns from experience."',
             size=16, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, bold=True)


def slide_05_pipeline(prs):
    """Slide 5 — Cognitive Pipeline."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "The Cognitive Pipeline")

    stages = ["Perceive", "Encode", "Store", "Consolidate", "Retrieve"]
    box_w = Inches(1.8)
    box_h = Inches(0.9)
    arrow_w = Inches(0.5)
    arrow_h = Inches(0.35)
    start_x = Inches(0.7)
    y = Inches(2.8)

    for i, stage in enumerate(stages):
        x = start_x + i * (box_w + arrow_w + Inches(0.1))
        add_rounded_rect(slide, x, y, box_w, box_h, stage, font_size=16)
        if i < len(stages) - 1:
            ax = x + box_w + Inches(0.05)
            ay = y + (box_h - arrow_h) / 2
            add_arrow(slide, ax, ay, arrow_w, arrow_h)

    loop_y = Inches(4.3)
    add_rounded_rect(slide, Inches(3.0), loop_y, Inches(1.6), Inches(0.7), "Dream", fill_color=DARK_GREEN_BG, border_color=ACCENT_GREEN, font_size=14)
    add_rounded_rect(slide, Inches(5.2), loop_y, Inches(1.6), Inches(0.7), "Abstract", fill_color=DARK_AMBER_BG, border_color=ACCENT_AMBER, font_size=14)
    add_rounded_rect(slide, Inches(7.4), loop_y, Inches(2.0), Inches(0.7), "Metacognition", fill_color=DARK_PURPLE_BG, border_color=ACCENT_PURPLE, font_size=14)
    add_rounded_rect(slide, Inches(10.0), loop_y, Inches(2.0), Inches(0.7), "Orchestrator", fill_color=DARK_RED_BG, border_color=ACCENT_RED, font_size=14)

    tf_main = add_textbox(slide, Inches(0.7), Inches(2.1), Inches(5), Inches(0.5))
    set_para(tf_main.paragraphs[0], "MAIN PIPELINE", size=12, bold=True, color=ACCENT_BLUE)
    tf_auto = add_textbox(slide, Inches(3.0), Inches(3.85), Inches(5), Inches(0.4))
    set_para(tf_auto.paragraphs[0], "AUTONOMOUS BACKGROUND AGENTS", size=12, bold=True, color=ACCENT_GREEN)

    tf_note = add_textbox(slide, Inches(0.8), Inches(5.5), Inches(11.7), Inches(0.8))
    set_para(tf_note.paragraphs[0],
             "Fully autonomous \u2014 runs 24/7 with zero user intervention. Event-driven pub/sub architecture.",
             size=16, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_06_agents(prs):
    """Slide 6 — 9 Agents."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "9 Cognitive Agents, Working Together")

    agents = [
        ("\U0001F441  Perception",     "Watches filesystem, terminal, clipboard \u2014 filters 80% noise"),
        ("\U0001F9E0  Encoding",       "Compresses raw events into structured memories with embeddings"),
        ("\U0001F4D6  Episoding",      "Groups events into coherent sessions with narrative synthesis"),
        ("\U0001F4A4  Consolidation",  "Sleep cycles: decay, merge, discover patterns"),
        ("\U0001F50D  Retrieval",      "Spread activation + multi-turn LLM synthesis"),
        ("\U0001F9D0  Metacognition",  "Self-monitoring, quality audits, feedback-driven adjustments"),
        ("\U0001F30C  Dreaming",       "Offline replay, cross-project linking, insight generation"),
        ("\U0001F4A1  Abstraction",    "Builds principles & axioms from recurring patterns"),
        ("\u2699\uFE0F  Orchestrator", "Health monitoring, self-testing, auto-recovery"),
    ]

    for col_idx, agent_slice in enumerate([agents[:5], agents[5:]]):
        x = Inches(0.6) + col_idx * Inches(6.2)
        tf = add_textbox(slide, x, Inches(1.7), Inches(6.0), Inches(5.5))
        tf.paragraphs[0].space_after = Pt(0)
        for i, (name, desc) in enumerate(agent_slice):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = name
            p.font.name = FONT_NAME
            p.font.size = Pt(16)
            p.font.bold = True
            p.font.color.rgb = ACCENT_BLUE
            p.space_before = Pt(7)
            p.space_after = Pt(1)
            dp = tf.add_paragraph()
            dp.text = desc
            dp.font.name = FONT_NAME
            dp.font.size = Pt(13)
            dp.font.color.rgb = LIGHT_GRAY
            dp.space_after = Pt(5)


def slide_07_innovation(prs):
    """Slide 7 — Technical Innovation."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "Neuroscience, Not Just Vectors")

    tf = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.2))
    tf.paragraphs[0].space_after = Pt(0)
    items = [
        ("Spread Activation", "Graph-based retrieval traverses 3 hops with 0.7 decay \u2014 finds connections vectors miss entirely"),
        ("Salience Decay", "Memories naturally fade (Ebbinghaus curve) \u2014 recently accessed memories resist decay"),
        ("Consolidation Cycles", "Merges related memories, discovers patterns, prunes noise \u2014 like what your brain does during sleep"),
        ("Association Graph", "Captures semantic relationships (caused_by, contradicts, reinforces) that cosine similarity cannot"),
    ]
    add_titled_bullets(tf, items, title_size=20, desc_size=15)


def slide_08_performance(prs):
    """Slide 8 — Performance + Cost Comparison."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "96% Retrieval Precision")

    # Hero number
    tf_big = add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3.5), Inches(2.0))
    set_para(tf_big.paragraphs[0], "96%", size=72, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
    p2 = tf_big.add_paragraph()
    p2.text = "average precision"
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT_NAME
    p2.font.size = Pt(16)
    p2.font.color.rgb = LIGHT_GRAY

    # Stats boxes
    stats = [("15", "Seed Memories"), ("5", "Benchmark Queries"), ("<500ms", "Query Latency"), ("3x", "Retrieval Methods")]
    for i, (num, label) in enumerate(stats):
        row, col = i // 2, i % 2
        add_stat_box(slide, Inches(5.0) + col * Inches(3.2), Inches(1.6) + row * Inches(1.5),
                     Inches(2.8), Inches(1.2), num, label)

    # Cost comparison section
    tf_label = add_textbox(slide, Inches(0.8), Inches(4.2), Inches(3), Inches(0.4))
    set_para(tf_label.paragraphs[0], "COST TO OPERATE", size=13, bold=True, color=ACCENT_GREEN)

    comparisons = [
        ("Per query:", "~$0 (local LLM)", "$0.01\u2013$0.10 (cloud APIs)"),
        ("10K memories:", "$0 infrastructure", "$200\u2013$500/mo (Pinecone, Weaviate)"),
        ("Data egress:", "$0 (air-gapped)", "$5K\u2013$50K/yr at enterprise scale"),
    ]
    tf_cost = add_textbox(slide, Inches(0.8), Inches(4.6), Inches(11.7), Inches(2.5))
    tf_cost.paragraphs[0].space_after = Pt(0)
    for i, (label, ours, theirs) in enumerate(comparisons):
        p = tf_cost.paragraphs[0] if i == 0 else tf_cost.add_paragraph()
        p.space_after = Pt(4)
        p.space_before = Pt(2)
        r_label = p.add_run()
        r_label.text = label + "  "
        r_label.font.name = FONT_NAME
        r_label.font.size = Pt(14)
        r_label.font.color.rgb = DIM_TEXT
        r_label.font.bold = True
        r_ours = p.add_run()
        r_ours.text = ours
        r_ours.font.name = FONT_NAME
        r_ours.font.size = Pt(14)
        r_ours.font.color.rgb = ACCENT_GREEN
        r_ours.font.bold = True
        r_vs = p.add_run()
        r_vs.text = "  vs  "
        r_vs.font.name = FONT_NAME
        r_vs.font.size = Pt(14)
        r_vs.font.color.rgb = DIM_TEXT
        r_them = p.add_run()
        r_them.text = theirs
        r_them.font.name = FONT_NAME
        r_them.font.size = Pt(14)
        r_them.font.color.rgb = X_RED


def slide_09_competitive(prs):
    """Slide 9 — NEW: Competitive Landscape."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "Where We Win")

    # Table header
    headers = ["Feature", "Mnemonic", "Mem0", "Zep", "LangChain"]
    col_widths = [Inches(2.8), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
    start_x = Inches(0.8)
    header_y = Inches(1.7)
    row_h = Inches(0.45)

    # Header row
    x = start_x
    for i, (header, w) in enumerate(zip(headers, col_widths)):
        tf = add_textbox(slide, x, header_y, w, row_h)
        color = ACCENT_BLUE if i == 0 else (ACCENT_GREEN if i == 1 else DIM_TEXT)
        set_para(tf.paragraphs[0], header, size=14, bold=True, color=color, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)
        x += w

    # Separator
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x, header_y + row_h, Inches(10.8), Pt(1.5))
    sep.fill.solid()
    sep.fill.fore_color.rgb = SHAPE_FILL
    sep.line.fill.background()

    # Data rows
    rows = [
        ("Local-first / air-gapped",  "\u2713", "\u2717", "\u2717", "Partial"),
        ("Cognitive agents",           "9",      "0",      "0",      "0"),
        ("Spread activation",          "\u2713", "\u2717", "\u2717", "\u2717"),
        ("Self-improving (feedback)",  "\u2713", "\u2717", "Partial", "\u2717"),
        ("Pattern discovery",          "\u2713", "\u2717", "\u2717", "\u2717"),
        ("MCP native (10 tools)",      "\u2713", "\u2717", "\u2717", "\u2717"),
        ("Dreaming / consolidation",   "\u2713", "\u2717", "\u2717", "\u2717"),
        ("Knowledge abstraction",      "\u2713", "\u2717", "\u2717", "\u2717"),
    ]

    for row_i, (feature, *vals) in enumerate(rows):
        y = header_y + row_h + Pt(4) + row_i * row_h
        x = start_x

        # Feature name
        tf = add_textbox(slide, x, y, col_widths[0], row_h)
        set_para(tf.paragraphs[0], feature, size=13, color=LIGHT_GRAY)
        x += col_widths[0]

        # Values
        for j, val in enumerate(vals):
            tf = add_textbox(slide, x, y, col_widths[j + 1], row_h)
            if val == "\u2713":
                c = CHECK_GREEN
            elif val == "\u2717":
                c = X_RED
            elif val == "Partial":
                c = ACCENT_AMBER
            else:
                c = WHITE
            is_bold = (j == 0)  # bold Mnemonic column
            set_para(tf.paragraphs[0], val, size=14, bold=is_bold, color=c, align=PP_ALIGN.CENTER)
            x += col_widths[j + 1]

    # Positioning statement
    tf_pos = add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.6))
    set_para(tf_pos.paragraphs[0],
             "Only solution combining neuroscience-grounded cognition with local-first privacy and MCP integration.",
             size=15, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, bold=True)


def slide_10_privacy(prs):
    """Slide 10 — Privacy & Local-First + Enterprise Premium."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "Your Data Never Leaves Your Machine")

    tf = add_textbox(slide, Inches(0.8), Inches(1.8), Inches(7.5), Inches(4.5))
    tf.paragraphs[0].space_after = Pt(0)
    items = [
        ("\U0001F512  Single Go binary + SQLite database", "No servers, no cloud dependencies, no attack surface"),
        ("\U0001F4BB  Local LLM via LM Studio", "Apple Silicon (~8B models) \u2014 zero API costs"),
        ("\u2708\uFE0F  Fully air-gapped operation", "Works offline, on planes, in secure environments"),
        ("\U0001F4BE  One-command backup", "Entire memory system in a single portable file"),
    ]
    add_titled_bullets(tf, items, title_size=18, desc_size=14)

    # Enterprise premium stats
    add_stat_box(slide, Inches(9.0), Inches(1.8), Inches(3.5), Inches(1.4), "43%", "of Fortune 500 testing\nlocal-only AI (Gartner)", num_color=ACCENT_AMBER, border=ACCENT_AMBER)
    add_stat_box(slide, Inches(9.0), Inches(3.5), Inches(3.5), Inches(1.4), "15\u201330%", "price premium for\nair-gapped solutions", num_color=ACCENT_GREEN, border=ACCENT_GREEN)
    add_stat_box(slide, Inches(9.0), Inches(5.2), Inches(3.5), Inches(1.4), "$0", "data egress costs\nvs $5K\u2013$50K/yr cloud", num_color=WHITE, border=ACCENT_BLUE)


def slide_11_mcp(prs):
    """Slide 11 — MCP Integration."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "Native AI Agent Integration")

    tools = [
        ("remember", "Store decisions, errors, insights"),
        ("recall", "Semantic search with synthesis"),
        ("forget", "Archive outdated memories"),
        ("status", "System health snapshot"),
        ("recall_project", "Project-scoped retrieval"),
        ("recall_timeline", "Time-range queries"),
        ("session_summary", "Current session recap"),
        ("get_patterns", "Discovered recurring themes"),
        ("get_insights", "Metacognition observations"),
        ("feedback", "Rate recall quality"),
    ]

    tf_label = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4))
    set_para(tf_label.paragraphs[0], "10 MCP TOOLS", size=14, bold=True, color=ACCENT_BLUE)

    tf_tools = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5))
    tf_tools.paragraphs[0].space_after = Pt(0)
    for i, (name, desc) in enumerate(tools):
        p = tf_tools.paragraphs[0] if i == 0 else tf_tools.add_paragraph()
        p.space_after = Pt(3)
        p.space_before = Pt(2)
        r_name = p.add_run()
        r_name.text = name
        r_name.font.name = MONO_FONT
        r_name.font.size = Pt(13)
        r_name.font.bold = True
        r_name.font.color.rgb = ACCENT_BLUE
        r_sep = p.add_run()
        r_sep.text = "  \u2014  "
        r_sep.font.name = FONT_NAME
        r_sep.font.size = Pt(13)
        r_sep.font.color.rgb = DIM_TEXT
        r_desc = p.add_run()
        r_desc.text = desc
        r_desc.font.name = FONT_NAME
        r_desc.font.size = Pt(13)
        r_desc.font.color.rgb = LIGHT_GRAY

    # Before/After boxes
    box_x = Inches(7.0)
    before_box = add_rounded_rect(slide, box_x, Inches(2.0), Inches(5.5), Inches(1.4), "", fill_color=DARK_RED_BG, border_color=ACCENT_RED)
    btf = before_box.text_frame
    btf.clear()
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    bp = btf.paragraphs[0]
    bp.alignment = PP_ALIGN.CENTER
    for text, sz, bld, clr in [("BEFORE\n", 11, True, ACCENT_RED),
                                ('"What did we decide about the database?"\n', 13, False, LIGHT_GRAY),
                                ("\U0001F937  No idea. Start over.", 13, False, DIM_TEXT)]:
        r = bp.add_run()
        r.text = text
        r.font.name = FONT_NAME
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr

    after_box = add_rounded_rect(slide, box_x, Inches(3.8), Inches(5.5), Inches(1.4), "", fill_color=DARK_GREEN_BG, border_color=ACCENT_GREEN)
    atf = after_box.text_frame
    atf.clear()
    atf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ap = atf.paragraphs[0]
    ap.alignment = PP_ALIGN.CENTER
    for text, sz, bld, clr in [("AFTER\n", 11, True, ACCENT_GREEN),
                                ('"What did we decide about the database?"\n', 13, False, LIGHT_GRAY),
                                ('\U0001F9E0  "We chose SQLite for simplicity. Here\'s why..."', 13, False, WHITE)]:
        r = ap.add_run()
        r.text = text
        r.font.name = FONT_NAME
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr

    # Ingest callout
    tf_ingest = add_textbox(slide, Inches(7.0), Inches(5.5), Inches(5.5), Inches(0.6))
    pi = tf_ingest.paragraphs[0]
    r = pi.add_run()
    r.text = "mnemonic ingest ."
    r.font.name = MONO_FONT
    r.font.size = Pt(15)
    r.font.bold = True
    r.font.color.rgb = ACCENT_BLUE
    r2 = pi.add_run()
    r2.text = "  \u2014  learns entire codebases autonomously"
    r2.font.name = FONT_NAME
    r2.font.size = Pt(14)
    r2.font.color.rgb = LIGHT_GRAY


def slide_12_pricing(prs):
    """Slide 12 — NEW: Revenue Model & Pricing."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "Three Paths to Revenue")

    tiers = [
        {
            "name": "Community",
            "price": "Free",
            "price_sub": "forever",
            "color": ACCENT_BLUE,
            "border": ACCENT_BLUE,
            "bg": RGBColor(0x1A, 0x35, 0x55),
            "features": [
                "Single user, local only",
                "10 MCP tools",
                "All 9 cognitive agents",
                "Unlimited memories",
                "Community support",
            ]
        },
        {
            "name": "Pro",
            "price": "$19/mo",
            "price_sub": "per user",
            "color": ACCENT_GREEN,
            "border": ACCENT_GREEN,
            "bg": DARK_GREEN_BG,
            "features": [
                "Everything in Community",
                "Team sync & sharing",
                "Cloud backup (encrypted)",
                "Priority encoding queue",
                "Analytics dashboard",
                "Email support",
            ]
        },
        {
            "name": "Enterprise",
            "price": "$49/user/mo",
            "price_sub": "or $15K\u2013$50K/yr site license",
            "color": ACCENT_AMBER,
            "border": ACCENT_AMBER,
            "bg": DARK_AMBER_BG,
            "features": [
                "Everything in Pro",
                "SSO / SAML / SCIM",
                "Audit logs & compliance",
                "Custom agents & plugins",
                "Air-gapped deployment",
                "Dedicated support + SLA",
            ]
        },
    ]

    for idx, tier in enumerate(tiers):
        x = Inches(0.6) + idx * Inches(4.2)
        w = Inches(3.9)

        # Tier box
        box = add_rounded_rect(slide, x, Inches(1.7), w, Inches(4.8), "", fill_color=tier["bg"], border_color=tier["border"])
        btf = box.text_frame
        btf.clear()
        btf.vertical_anchor = MSO_ANCHOR.TOP

        # Tier name
        p = btf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(12)
        r = p.add_run()
        r.text = tier["name"]
        r.font.name = FONT_NAME
        r.font.size = Pt(16)
        r.font.bold = True
        r.font.color.rgb = tier["color"]

        # Price
        pp = btf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        pp.space_before = Pt(6)
        rp = pp.add_run()
        rp.text = tier["price"]
        rp.font.name = FONT_NAME
        rp.font.size = Pt(28)
        rp.font.bold = True
        rp.font.color.rgb = WHITE

        # Price subtitle
        ps = btf.add_paragraph()
        ps.alignment = PP_ALIGN.CENTER
        ps.space_after = Pt(12)
        rs = ps.add_run()
        rs.text = tier["price_sub"]
        rs.font.name = FONT_NAME
        rs.font.size = Pt(11)
        rs.font.color.rgb = DIM_TEXT

        # Features
        for feat in tier["features"]:
            fp = btf.add_paragraph()
            fp.alignment = PP_ALIGN.LEFT
            fp.space_before = Pt(2)
            fp.space_after = Pt(2)
            fr = fp.add_run()
            fr.text = "  \u2713  " + feat
            fr.font.name = FONT_NAME
            fr.font.size = Pt(12)
            fr.font.color.rgb = LIGHT_GRAY

    # Bottom callout
    tf_market = add_textbox(slide, Inches(0.8), Inches(6.7), Inches(11.7), Inches(0.5))
    set_para(tf_market.paragraphs[0],
             "Plugin marketplace: 20% rev share on third-party agent plugins  \u2022  Benchmarks: Copilot $19/mo  \u2022  Cursor $20/mo  \u2022  Linear $8/mo",
             size=12, color=DIM_TEXT, align=PP_ALIGN.CENTER)


def slide_13_projections(prs):
    """Slide 13 — NEW: Financial Projections (3-Year)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "The Path to $12M ARR")

    # Projection table
    tf_label = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4))
    set_para(tf_label.paragraphs[0], "3-YEAR FINANCIAL PROJECTIONS", size=13, bold=True, color=ACCENT_BLUE)

    headers = ["", "Year 1", "Year 2", "Year 3"]
    col_ws = [Inches(2.2), Inches(1.8), Inches(1.8), Inches(1.8)]
    start_x = Inches(0.8)
    start_y = Inches(2.0)
    rh = Inches(0.45)

    # Headers
    x = start_x
    for h, w in zip(headers, col_ws):
        tf = add_textbox(slide, x, start_y, w, rh)
        set_para(tf.paragraphs[0], h, size=13, bold=True, color=ACCENT_BLUE,
                 align=PP_ALIGN.CENTER if h else PP_ALIGN.LEFT)
        x += w

    rows = [
        ("Free Users",     "5,000",   "25,000",  "80,000"),
        ("Pro Users",      "500",     "2,500",   "8,000"),
        ("Enterprise Orgs","10",      "30",      "80"),
        ("MRR",            "$25K",    "$150K",   "$1M"),
        ("ARR",            "$300K",   "$1.8M",   "$12M"),
        ("Gross Margin",   "90%",     "91%",     "92%"),
    ]

    for ri, (label, *vals) in enumerate(rows):
        y = start_y + (ri + 1) * rh + Pt(4)
        x = start_x

        # Highlight ARR row
        is_arr = (label == "ARR")
        label_color = ACCENT_AMBER if is_arr else LIGHT_GRAY
        val_color = ACCENT_AMBER if is_arr else WHITE

        tf = add_textbox(slide, x, y, col_ws[0], rh)
        set_para(tf.paragraphs[0], label, size=13, bold=is_arr, color=label_color)
        x += col_ws[0]

        for j, val in enumerate(vals):
            tf = add_textbox(slide, x, y, col_ws[j + 1], rh)
            set_para(tf.paragraphs[0], val, size=14, bold=is_arr, color=val_color, align=PP_ALIGN.CENTER)
            x += col_ws[j + 1]

    # Key assumptions box on right
    tf_al = add_textbox(slide, Inches(8.2), Inches(1.6), Inches(4.5), Inches(0.4))
    set_para(tf_al.paragraphs[0], "KEY ASSUMPTIONS", size=13, bold=True, color=ACCENT_GREEN)

    assume_box = add_rounded_rect(slide, Inches(8.2), Inches(2.0), Inches(4.5), Inches(3.8), "", fill_color=DARK_GREEN_BG, border_color=ACCENT_GREEN)
    atf = assume_box.text_frame
    atf.clear()
    atf.vertical_anchor = MSO_ANCHOR.TOP

    assumptions = [
        "10% free \u2192 paid conversion",
        "$19/mo Pro tier",
        "$49/user/mo Enterprise (avg 15 seats)",
        "5% monthly churn (Pro)",
        "1% monthly churn (Enterprise)",
        "120% net revenue retention",
        "Community-led growth (no paid ads Y1)",
        "$50 blended CAC (organic + content)",
    ]

    for i, text in enumerate(assumptions):
        p = atf.paragraphs[0] if i == 0 else atf.add_paragraph()
        p.text = "  \u2022  " + text
        p.font.name = FONT_NAME
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        if i == 0:
            p.space_before = Pt(10)

    # Visual ARR growth bar at bottom
    tf_growth = add_textbox(slide, Inches(0.8), Inches(5.5), Inches(7.0), Inches(0.4))
    set_para(tf_growth.paragraphs[0], "ARR GROWTH TRAJECTORY", size=12, bold=True, color=ACCENT_AMBER)

    # Three bars showing relative growth
    bar_data = [("Y1: $300K", 1.0), ("Y2: $1.8M", 3.0), ("Y3: $12M", 8.0)]
    max_w = Inches(6.5)
    for i, (label, scale) in enumerate(bar_data):
        by = Inches(5.9) + i * Inches(0.42)
        bw = max_w * (scale / 8.0)
        bar = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), by, bw, Inches(0.32))
        bar.fill.solid()
        colors = [ACCENT_BLUE, ACCENT_GREEN, ACCENT_AMBER]
        bar.fill.fore_color.rgb = colors[i]
        bar.line.fill.background()

        bar_tf = bar.text_frame
        bar_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        bar_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        r = bar_tf.paragraphs[0].add_run()
        r.text = "  " + label
        r.font.name = FONT_NAME
        r.font.size = Pt(11)
        r.font.bold = True
        r.font.color.rgb = WHITE


def slide_14_unit_economics(prs):
    """Slide 14 — NEW: Unit Economics."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "The Numbers That Matter")

    # Top row: 4 big metric boxes
    metrics_top = [
        ("$50", "Customer\nAcquisition Cost", "community-led, organic", ACCENT_BLUE, ACCENT_BLUE),
        ("$684", "LTV (Pro)", "36-mo avg \u00d7 $19/mo", ACCENT_GREEN, ACCENT_GREEN),
        ("$29,400", "LTV (Enterprise)", "5-yr \u00d7 $49 \u00d7 15 seats", ACCENT_AMBER, ACCENT_AMBER),
        ("13.7:1", "LTV:CAC (Pro)", "588:1 for Enterprise", WHITE, ACCENT_BLUE),
    ]

    for i, (num, label, sub, num_clr, border) in enumerate(metrics_top):
        x = Inches(0.6) + i * Inches(3.15)
        box = add_rounded_rect(slide, x, Inches(1.7), Inches(2.9), Inches(2.0), "", border_color=border)
        btf = box.text_frame
        btf.clear()
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = btf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.name = FONT_NAME
        r.font.size = Pt(32)
        r.font.bold = True
        r.font.color.rgb = num_clr

        p2 = btf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = label
        r2.font.name = FONT_NAME
        r2.font.size = Pt(13)
        r2.font.bold = True
        r2.font.color.rgb = LIGHT_GRAY

        p3 = btf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(2)
        r3 = p3.add_run()
        r3.text = sub
        r3.font.name = FONT_NAME
        r3.font.size = Pt(10)
        r3.font.color.rgb = DIM_TEXT

    # Bottom row: 3 operational metrics
    metrics_bottom = [
        ("92%", "Gross Margin", "No cloud infra for free/pro\u2014user's own hardware", ACCENT_GREEN, ACCENT_GREEN),
        ("<3 mo", "Payback Period (Pro)", "Enterprise: <1 month", ACCENT_BLUE, ACCENT_BLUE),
        ("Month 14", "Break-Even", "At current estimated burn rate", ACCENT_AMBER, ACCENT_AMBER),
    ]

    for i, (num, label, sub, num_clr, border) in enumerate(metrics_bottom):
        x = Inches(0.6) + i * Inches(4.2)
        box = add_rounded_rect(slide, x, Inches(4.2), Inches(3.9), Inches(1.8), "", border_color=border)
        btf = box.text_frame
        btf.clear()
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE

        p = btf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = num
        r.font.name = FONT_NAME
        r.font.size = Pt(30)
        r.font.bold = True
        r.font.color.rgb = num_clr

        p2 = btf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)
        r2 = p2.add_run()
        r2.text = label
        r2.font.name = FONT_NAME
        r2.font.size = Pt(13)
        r2.font.bold = True
        r2.font.color.rgb = LIGHT_GRAY

        p3 = btf.add_paragraph()
        p3.alignment = PP_ALIGN.CENTER
        p3.space_before = Pt(2)
        r3 = p3.add_run()
        r3.text = sub
        r3.font.name = FONT_NAME
        r3.font.size = Pt(10)
        r3.font.color.rgb = DIM_TEXT

    # Bottom callout
    tf_bot = add_textbox(slide, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5))
    set_para(tf_bot.paragraphs[0],
             "Local-first = near-zero marginal cost. Every user runs on their own hardware. We ship software, not servers.",
             size=13, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, bold=True)


def slide_15_gtm(prs):
    """Slide 15 — NEW: Go-to-Market Strategy."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "How We Get There")

    phases = [
        {
            "title": "PHASE 1: TRACTION",
            "period": "Months 1\u20136",
            "color": ACCENT_BLUE,
            "border": ACCENT_BLUE,
            "bg": RGBColor(0x1A, 0x35, 0x55),
            "items": [
                "Open-source launch",
                "GitHub / HN / Product Hunt",
                "MCP ecosystem (16K+ servers)",
                "Claude Code community",
            ],
            "target": "5,000 installs\n500 daily active users",
        },
        {
            "title": "PHASE 2: MONETIZE",
            "period": "Months 6\u201312",
            "color": ACCENT_GREEN,
            "border": ACCENT_GREEN,
            "bg": DARK_GREEN_BG,
            "items": [
                "Launch Pro tier ($19/mo)",
                "Team sync feature",
                "Plugin marketplace (20+ agents)",
                "Content marketing + docs",
            ],
            "target": "$25K MRR\n10 enterprise pilots",
        },
        {
            "title": "PHASE 3: SCALE",
            "period": "Months 12\u201324",
            "color": ACCENT_AMBER,
            "border": ACCENT_AMBER,
            "bg": DARK_AMBER_BG,
            "items": [
                "Enterprise sales motion",
                "SOC 2 / HIPAA compliance",
                "Channel partnerships",
                "International expansion",
            ],
            "target": "$150K MRR\n30 enterprise contracts",
        },
    ]

    for idx, phase in enumerate(phases):
        x = Inches(0.6) + idx * Inches(4.2)
        w = Inches(3.9)

        box = add_rounded_rect(slide, x, Inches(1.7), w, Inches(4.0), "", fill_color=phase["bg"], border_color=phase["border"])
        btf = box.text_frame
        btf.clear()
        btf.vertical_anchor = MSO_ANCHOR.TOP

        # Title
        p = btf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        p.space_before = Pt(10)
        r = p.add_run()
        r.text = phase["title"]
        r.font.name = FONT_NAME
        r.font.size = Pt(14)
        r.font.bold = True
        r.font.color.rgb = phase["color"]

        # Period
        pp = btf.add_paragraph()
        pp.alignment = PP_ALIGN.CENTER
        pp.space_after = Pt(10)
        rp = pp.add_run()
        rp.text = phase["period"]
        rp.font.name = FONT_NAME
        rp.font.size = Pt(11)
        rp.font.color.rgb = DIM_TEXT

        # Items
        for item in phase["items"]:
            ip = btf.add_paragraph()
            ip.alignment = PP_ALIGN.LEFT
            ip.space_before = Pt(2)
            ip.space_after = Pt(2)
            ir = ip.add_run()
            ir.text = "  \u2022  " + item
            ir.font.name = FONT_NAME
            ir.font.size = Pt(12)
            ir.font.color.rgb = LIGHT_GRAY

        # Target
        tp = btf.add_paragraph()
        tp.alignment = PP_ALIGN.CENTER
        tp.space_before = Pt(12)
        tr = tp.add_run()
        tr.text = "\u2192 " + phase["target"]
        tr.font.name = FONT_NAME
        tr.font.size = Pt(12)
        tr.font.bold = True
        tr.font.color.rgb = WHITE

    # Arrows between phases
    for i in range(2):
        ax = Inches(4.5) + i * Inches(4.2)
        ay = Inches(3.5)
        add_arrow(slide, ax, ay, Inches(0.3), Inches(0.25))

    # Bottom: growth flywheel
    tf_fly = add_textbox(slide, Inches(0.8), Inches(6.0), Inches(11.7), Inches(0.8))
    set_para(tf_fly.paragraphs[0],
             "Open Source \u2192 Community \u2192 Adoption \u2192 Feedback \u2192 Product Improvement \u2192 Enterprise Demand",
             size=14, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, bold=True)
    p2 = tf_fly.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.font.name = FONT_NAME
    p2.font.size = Pt(12)
    p2.font.color.rgb = DIM_TEXT
    p2.text = "The open-source flywheel: users become advocates, advocates become customers, customers become champions."


def slide_16_usecases(prs):
    """Slide 16 — Use Cases + Revenue Per Segment."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "Who Pays For This?")

    segments = [
        {
            "icon": "\U0001F468\u200D\U0001F4BB",
            "name": "AI-Assisted Developers",
            "desc": "Persistent context across coding sessions",
            "market": "12M+ developers using AI tools",
            "revenue": "$2B addressable",
            "color": ACCENT_BLUE,
        },
        {
            "icon": "\U0001F3E2",
            "name": "Enterprise Teams",
            "desc": "Institutional knowledge that never gets lost",
            "market": "53% of enterprises use AI agents",
            "revenue": "$500K+ avg deal size",
            "color": ACCENT_GREEN,
        },
        {
            "icon": "\U0001F512",
            "name": "Security / Government",
            "desc": "Air-gapped, compliance-ready memory",
            "market": "43% Fortune 500 want local-only AI",
            "revenue": "15\u201330% price premium",
            "color": ACCENT_AMBER,
        },
        {
            "icon": "\U0001F52C",
            "name": "Research Labs",
            "desc": "Cross-project pattern discovery",
            "market": "Growing AI research spend",
            "revenue": "$10K+/yr per lab",
            "color": ACCENT_PURPLE,
        },
    ]

    for i, seg in enumerate(segments):
        row = i // 2
        col = i % 2
        x = Inches(0.6) + col * Inches(6.3)
        y = Inches(1.7) + row * Inches(2.7)
        w = Inches(6.0)
        h = Inches(2.4)

        box = add_rounded_rect(slide, x, y, w, h, "", border_color=seg["color"])
        btf = box.text_frame
        btf.clear()
        btf.vertical_anchor = MSO_ANCHOR.TOP

        # Name
        p = btf.paragraphs[0]
        p.space_before = Pt(10)
        r = p.add_run()
        r.text = seg["icon"] + "  " + seg["name"]
        r.font.name = FONT_NAME
        r.font.size = Pt(18)
        r.font.bold = True
        r.font.color.rgb = seg["color"]

        # Description
        dp = btf.add_paragraph()
        dp.space_before = Pt(4)
        dr = dp.add_run()
        dr.text = seg["desc"]
        dr.font.name = FONT_NAME
        dr.font.size = Pt(13)
        dr.font.color.rgb = LIGHT_GRAY

        # Market size
        mp = btf.add_paragraph()
        mp.space_before = Pt(8)
        mr = mp.add_run()
        mr.text = "\U0001F4CA  " + seg["market"]
        mr.font.name = FONT_NAME
        mr.font.size = Pt(13)
        mr.font.color.rgb = WHITE

        # Revenue
        rp = btf.add_paragraph()
        rp.space_before = Pt(3)
        rr = rp.add_run()
        rr.text = "\U0001F4B0  " + seg["revenue"]
        rr.font.name = FONT_NAME
        rr.font.size = Pt(13)
        rr.font.bold = True
        rr.font.color.rgb = seg["color"]


def slide_17_architecture(prs):
    """Slide 17 — Architecture (tech details)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_NAVY)
    add_slide_title(slide, "Built to Last")

    # Left: tech stack
    tf_label = add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4))
    set_para(tf_label.paragraphs[0], "TECH STACK", size=13, bold=True, color=ACCENT_BLUE)

    tf_specs = add_textbox(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(4.5))
    tf_specs.paragraphs[0].space_after = Pt(0)
    specs = [
        ("Go 1.23", "Single binary, zero runtime deps, CGO-enabled"),
        ("SQLite + WAL", "Sub-ms reads, concurrent access, FTS5 search"),
        ("Event Bus", "Loosely coupled pub/sub \u2014 add agents without touching core"),
        ("Pluggable", "Swap LLM providers, stores, watchers independently"),
        ("Budget-Constrained", "Max 100 memories/cycle, 5 merges, 3 tool calls"),
    ]
    add_titled_bullets(tf_specs, specs, title_size=16, desc_size=13)

    # Right: data model
    tf_label2 = add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5), Inches(0.4))
    set_para(tf_label2.paragraphs[0], "DATA MODEL", size=13, bold=True, color=ACCENT_BLUE)

    schema = [
        "memories \u2014 encoded knowledge with embeddings",
        "associations \u2014 weighted graph edges",
        "episodes \u2014 temporal session clusters",
        "patterns \u2014 recurring themes with evidence",
        "abstractions \u2014 3-level hierarchy",
        "raw_memories \u2014 unprocessed observation queue",
    ]
    tf_schema = add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(3.0))
    tf_schema.paragraphs[0].space_after = Pt(0)
    for i, item in enumerate(schema):
        p = tf_schema.paragraphs[0] if i == 0 else tf_schema.add_paragraph()
        p.text = "\u2022  " + item
        p.font.name = FONT_NAME
        p.font.size = Pt(13)
        p.font.color.rgb = LIGHT_GRAY
        p.space_after = Pt(6)

    insight = add_rounded_rect(slide, Inches(7.0), Inches(5.0), Inches(5.5), Inches(1.0), "", font_size=14)
    itf = insight.text_frame
    itf.clear()
    itf.vertical_anchor = MSO_ANCHOR.MIDDLE
    ip = itf.paragraphs[0]
    ip.alignment = PP_ALIGN.CENTER
    r1 = ip.add_run()
    r1.text = "The cognitive model is durable. "
    r1.font.name = FONT_NAME
    r1.font.size = Pt(14)
    r1.font.bold = True
    r1.font.color.rgb = WHITE
    r2 = ip.add_run()
    r2.text = "Implementations are swappable."
    r2.font.name = FONT_NAME
    r2.font.size = Pt(14)
    r2.font.color.rgb = ACCENT_BLUE


def slide_18_cta(prs):
    """Slide 18 — Call to Action with financial targets."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_DARK)

    tf = add_textbox(slide, Inches(1), Inches(0.8), Inches(11.3), Inches(1.0))
    set_para(tf.paragraphs[0], "Let's Build the Future of AI Memory", size=38, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.2), Inches(1.8), Inches(3), Pt(3))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Three milestone boxes
    milestones = [
        ("12 Months", "$300K ARR", "5,000 users  \u2022  10 enterprise pilots", ACCENT_BLUE),
        ("24 Months", "$1.8M ARR", "25,000 users  \u2022  Path to profitability", ACCENT_GREEN),
        ("36 Months", "$12M ARR", "80,000 users  \u2022  80 enterprise contracts", ACCENT_AMBER),
    ]

    for i, (period, arr, detail, color) in enumerate(milestones):
        x = Inches(0.6) + i * Inches(4.2)
        box = add_rounded_rect(slide, x, Inches(2.2), Inches(3.9), Inches(1.8), "", border_color=color)
        btf = box.text_frame
        btf.clear()
        btf.vertical_anchor = MSO_ANCHOR.MIDDLE

        pp = btf.paragraphs[0]
        pp.alignment = PP_ALIGN.CENTER
        rr = pp.add_run()
        rr.text = period
        rr.font.name = FONT_NAME
        rr.font.size = Pt(13)
        rr.font.bold = True
        rr.font.color.rgb = color

        ap = btf.add_paragraph()
        ap.alignment = PP_ALIGN.CENTER
        ap.space_before = Pt(4)
        ar = ap.add_run()
        ar.text = arr
        ar.font.name = FONT_NAME
        ar.font.size = Pt(30)
        ar.font.bold = True
        ar.font.color.rgb = WHITE

        dp = btf.add_paragraph()
        dp.alignment = PP_ALIGN.CENTER
        dp.space_before = Pt(4)
        dr = dp.add_run()
        dr.text = detail
        dr.font.name = FONT_NAME
        dr.font.size = Pt(11)
        dr.font.color.rgb = LIGHT_GRAY

    # Key metrics strip
    tf_metrics = add_textbox(slide, Inches(0.8), Inches(4.3), Inches(11.7), Inches(0.6))
    set_para(tf_metrics.paragraphs[0],
             "92% gross margins  \u2022  13:1 LTV:CAC  \u2022  Community-led growth  \u2022  Zero cloud infrastructure costs",
             size=16, bold=True, color=ACCENT_GREEN, align=PP_ALIGN.CENTER)

    # What's already built
    tf_built = add_textbox(slide, Inches(1.5), Inches(5.0), Inches(10.3), Inches(1.0))
    tf_built.paragraphs[0].space_after = Pt(0)
    built_items = [
        "\u2705  v1.0 is live \u2014 9 agents, 10 MCP tools, 96% precision, benchmark-validated",
        "\u2705  The hard part is done \u2014 neuroscience-grounded architecture, working product",
        "\U0001F680  The moat is real \u2014 cognitive model + local-first + MCP ecosystem lock-in",
    ]
    for i, text in enumerate(built_items):
        p = tf_built.paragraphs[0] if i == 0 else tf_built.add_paragraph()
        p.text = text
        p.font.name = FONT_NAME
        p.font.size = Pt(15)
        p.font.color.rgb = WHITE
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.CENTER

    # Closing
    tf_close = add_textbox(slide, Inches(1.5), Inches(6.3), Inches(10.3), Inches(0.7))
    set_para(tf_close.paragraphs[0],
             '"Every AI agent deserves a memory. Let\'s give it to them."',
             size=20, color=ACCENT_BLUE, align=PP_ALIGN.CENTER, bold=True)


# ═══════════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════════

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)        #  1: Title
    slide_02_problem(prs)      #  2: Problem + market stats
    slide_03_market(prs)       #  3: TAM/SAM/SOM
    slide_04_solution(prs)     #  4: Meet Mnemonic
    slide_05_pipeline(prs)     #  5: Cognitive Pipeline
    slide_06_agents(prs)       #  6: 9 Agents
    slide_07_innovation(prs)   #  7: Neuroscience innovation
    slide_08_performance(prs)  #  8: 96% precision + cost comparison
    slide_09_competitive(prs)  #  9: Competitive landscape
    slide_10_privacy(prs)      # 10: Privacy + enterprise premium
    slide_11_mcp(prs)          # 11: MCP integration
    slide_12_pricing(prs)      # 12: Revenue model / pricing tiers
    slide_13_projections(prs)  # 13: 3-year financial projections
    slide_14_unit_economics(prs) # 14: Unit economics
    slide_15_gtm(prs)          # 15: Go-to-market strategy
    slide_16_usecases(prs)     # 16: Use cases + revenue per segment
    slide_17_architecture(prs) # 17: Architecture
    slide_18_cta(prs)          # 18: Call to action with targets

    prs.save(OUTPUT_PATH)
    print(f"\nPresentation saved to: {OUTPUT_PATH}")
    print(f"  Slides: {len(prs.slides)}")
    print(f"  Format: 16:9 widescreen")
    print(f"  Theme:  Dark navy + electric blue accents")
    print(f"  New:    6 business/financial slides added")


if __name__ == "__main__":
    main()
